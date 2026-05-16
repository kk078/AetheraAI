"""
Aethera AI — Office Document Extractors
Extracts text and structure from DOCX, PPTX, and XLSX files.
"""
import logging
from typing import Dict, Any, List

logger = logging.getLogger("aethera.pipeline.office_docs")


def extract_docx(file_path: str) -> Dict[str, Any]:
    """
    Extract text and structure from a DOCX file.

    Returns:
        Dict with text, tables, metadata
    """
    result = {"text": "", "tables": [], "metadata": {}, "success": False}

    try:
        from docx import Document
    except ImportError:
        result["error"] = "python-docx not installed"
        return result

    try:
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                result["tables"].append({"rows": table_data, "headers": table_data[0] if table_data else []})

        result["text"] = "\n\n".join(paragraphs)

        # Extract metadata
        core_props = doc.core_properties
        result["metadata"] = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
        }
        result["success"] = True

    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        result["error"] = str(e)

    return result


def extract_pptx(file_path: str) -> Dict[str, Any]:
    """
    Extract text from a PPTX file.

    Returns:
        Dict with text, tables, metadata
    """
    result = {"text": "", "tables": [], "metadata": {}, "success": False}

    try:
        from pptx import Presentation
    except ImportError:
        result["error"] = "python-pptx not installed"
        return result

    try:
        prs = Presentation(file_path)
        slides_text = []

        for slide in prs.slides:
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        result["tables"].append({"rows": table_data, "headers": table_data[0] if table_data else []})
            if slide_parts:
                slides_text.append("\n".join(slide_parts))

        result["text"] = "\n\n---\n\n".join(slides_text)
        result["metadata"]["slide_count"] = len(prs.slides)
        result["success"] = True

    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        result["error"] = str(e)

    return result


def extract_xlsx(file_path: str) -> Dict[str, Any]:
    """
    Extract text representation from an XLSX file.

    Returns:
        Dict with text (tab-separated representation), tables, metadata
    """
    result = {"text": "", "tables": [], "metadata": {}, "success": False}

    try:
        from openpyxl import load_workbook
    except ImportError:
        result["error"] = "openpyxl not installed"
        return result

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data = []
            for row in ws.iter_rows(max_row=500, values_only=True):
                row_strs = [str(cell) if cell is not None else "" for cell in row]
                rows_data.append(row_strs)
                if len(rows_data) > 500:
                    break

            if rows_data:
                result["tables"].append({
                    "sheet": sheet_name,
                    "headers": rows_data[0] if rows_data else [],
                    "rows": rows_data[:100],
                    "total_rows": len(rows_data),
                })
                # Text representation for indexing
                for row in rows_data[:50]:
                    text_parts.append("\t".join(row))

        result["text"] = "\n".join(text_parts)
        result["metadata"]["sheet_names"] = wb.sheetnames
        result["success"] = True
        wb.close()

    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        result["error"] = str(e)

    return result