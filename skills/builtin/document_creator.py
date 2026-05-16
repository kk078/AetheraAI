"""
Aethera AI - Document Creator Skill

Generate DOCX, PDF, XLSX, PPTX, and Markdown documents.
Uses python-docx, openpyxl, python-pptx when available;
falls back to structured JSON when libraries are absent.
"""

import io
import os
import re
from datetime import datetime
from typing import Optional

from skills.skill_base import AetheraSkill, SkillResult, skill

# Optional binary format libraries
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from fpdf import FPDF
    HAS_FPDF = True
except ImportError:
    HAS_FPDF = False


@skill(name="document_creator", category="general")
class DocumentCreatorSkill(AetheraSkill):
    """
    Create documents in various formats: DOCX, PDF, XLSX, PPTX, MD.
    Generates actual binary files when libraries are available.
    """

    @property
    def name(self) -> str:
        return "document_creator"

    @property
    def description(self) -> str:
        return "Generate documents in DOCX, PDF, XLSX, PPTX, or Markdown format"

    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "format": {
                    "type": "string",
                    "enum": ["docx", "pdf", "xlsx", "pptx", "md"],
                    "description": "Output format"
                },
                "content": {
                    "type": "string",
                    "description": "Document content (Markdown supported for all formats)"
                },
                "title": {
                    "type": "string",
                    "description": "Document title"
                },
                "template": {
                    "type": "string",
                    "description": "Template name to use"
                },
                "metadata": {
                    "type": "object",
                    "description": "Additional metadata (author, subject, etc.)"
                },
                "output_path": {
                    "type": "string",
                    "description": "File path to save the document (if omitted, returns base64)"
                }
            },
            "required": ["format", "content"]
        }

    @property
    def examples(self) -> list:
        return [
            {"input": {"format": "md", "content": "# Report\n\nContent here"}},
            {"input": {"format": "docx", "title": "Memo", "content": "To: Team\n\nMeeting notes..."}},
            {"input": {"format": "xlsx", "title": "Data", "content": "Name,Age,City\nAlice,30,NYC\nBob,25,LA"}},
        ]

    @property
    def cache_ttl(self) -> int:
        return 0

    async def execute(self, **kwargs) -> SkillResult:
        fmt = kwargs.get("format", "md").lower()
        content = kwargs.get("content", "")
        title = kwargs.get("title", "Document")
        metadata = kwargs.get("metadata", {})
        output_path = kwargs.get("output_path", "")

        if not content:
            return SkillResult(success=False, error="Content is required")

        try:
            if fmt == "md":
                result = self._create_markdown(content, title, metadata)
            elif fmt == "docx":
                result = self._create_docx(content, title, metadata, output_path)
            elif fmt == "xlsx":
                result = self._create_xlsx(content, title, metadata, output_path)
            elif fmt == "pptx":
                result = self._create_pptx(content, title, metadata, output_path)
            elif fmt == "pdf":
                result = self._create_pdf(content, title, metadata, output_path)
            else:
                return SkillResult(success=False, error=f"Unsupported format: {fmt}")
            return result
        except Exception as e:
            return SkillResult(success=False, error=str(e))

    def _save_or_b64(self, buffer: io.BytesIO, filename: str, output_path: str) -> dict:
        """Save buffer to file path, or return base64-encoded content."""
        import base64
        data = buffer.getvalue()

        if output_path:
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            with open(output_path, "wb") as f:
                f.write(data)
            return {"filename": filename, "path": output_path, "size_bytes": len(data)}

        return {
            "filename": filename,
            "content_base64": base64.b64encode(data).decode("ascii"),
            "size_bytes": len(data)
        }

    # ------------------------------------------------------------------
    # Markdown
    # ------------------------------------------------------------------

    def _create_markdown(self, content: str, title: str, metadata: dict) -> SkillResult:
        """Create Markdown document."""
        output = f"# {title}\n\n"
        if metadata.get("author"):
            output += f"*By {metadata['author']}*\n\n"
        if metadata.get("date"):
            output += f"*{metadata['date']}*\n\n"
        elif metadata.get("include_date", True):
            output += f"*{datetime.now().strftime('%B %d, %Y')}*\n\n"
        output += content

        return SkillResult(
            success=True,
            data={
                "format": "markdown",
                "content": output,
                "filename": self._safe_filename(title, ".md")
            }
        )

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _create_docx(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Create Word document using python-docx or structured fallback."""
        if HAS_DOCX:
            return self._create_docx_binary(content, title, metadata, output_path)

        # Fallback: return structured JSON
        doc_structure = self._parse_markdown_to_structure(content, title, metadata)
        return SkillResult(
            success=True,
            data={
                "format": "docx",
                "structure": doc_structure,
                "filename": self._safe_filename(title, ".docx"),
                "note": "python-docx not installed — install with: pip install python-docx"
            }
        )

    def _create_docx_binary(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Generate actual DOCX using python-docx."""
        doc = Document()

        # Set metadata
        doc.core_properties.title = title
        doc.core_properties.author = metadata.get("author", "Aethera AI")
        doc.core_properties.created = datetime.now()

        # Parse markdown content into paragraphs
        lines = content.split("\n")
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped.startswith("1. ") or stripped.startswith("2. "):
                doc.add_paragraph(stripped[3:], style="List Number")
            elif stripped.startswith("**") and stripped.endswith("**"):
                p = doc.add_paragraph()
                run = p.add_run(stripped[2:-2])
                run.bold = True
            else:
                doc.add_paragraph(stripped)

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        file_info = self._save_or_b64(buffer, self._safe_filename(title, ".docx"), output_path)
        file_info["format"] = "docx"
        return SkillResult(success=True, data=file_info)

    # ------------------------------------------------------------------
    # XLSX
    # ------------------------------------------------------------------

    def _create_xlsx(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Create Excel spreadsheet."""
        if HAS_OPENPYXL:
            return self._create_xlsx_binary(content, title, metadata, output_path)

        # Fallback: return structured data
        rows = self._parse_csv_content(content)
        return SkillResult(
            success=True,
            data={
                "format": "xlsx",
                "title": title,
                "rows": rows,
                "filename": self._safe_filename(title, ".xlsx"),
                "note": "openpyxl not installed — install with: pip install openpyxl"
            }
        )

    def _create_xlsx_binary(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Generate actual XLSX using openpyxl."""
        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name limit

        # Parse content as CSV-like data
        rows = self._parse_csv_content(content)

        header_font = Font(bold=True, size=11)
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font_white = Font(bold=True, size=11, color="FFFFFF")
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin")
        )

        for row_idx, row in enumerate(rows, 1):
            for col_idx, cell_value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                cell.border = thin_border
                if row_idx == 1:
                    cell.font = header_font_white
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")

        # Auto-width columns
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)

        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)

        file_info = self._save_or_b64(buffer, self._safe_filename(title, ".xlsx"), output_path)
        file_info["format"] = "xlsx"
        file_info["rows"] = len(rows)
        file_info["columns"] = len(rows[0]) if rows else 0
        return SkillResult(success=True, data=file_info)

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _create_pptx(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Create PowerPoint presentation."""
        if HAS_PPTX:
            return self._create_pptx_binary(content, title, metadata, output_path)

        # Fallback: return structured data
        slides = self._parse_slides(content, title)
        return SkillResult(
            success=True,
            data={
                "format": "pptx",
                "title": title,
                "slides": slides,
                "filename": self._safe_filename(title, ".pptx"),
                "note": "python-pptx not installed — install with: pip install python-pptx"
            }
        )

    def _create_pptx_binary(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Generate actual PPTX using python-pptx."""
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = metadata.get("author", "Aethera AI")

        # Content slides (split by --- or ## headings)
        content_sections = re.split(r"\n---\n|\n##\s", content)
        for section in content_sections:
            section = section.strip()
            if not section:
                continue

            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            lines = section.split("\n")
            # First line as title, rest as body
            slide_title = lines[0].lstrip("#").strip() if lines else "Slide"
            slide.shapes.title.text = slide_title[:50]

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                # Strip markdown bullet markers
                line = line.lstrip("-*0123456789). ")
                if line:
                    p = tf.add_paragraph()
                    p.text = line[:200]
                    p.level = 0

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        file_info = self._save_or_b64(buffer, self._safe_filename(title, ".pptx"), output_path)
        file_info["format"] = "pptx"
        file_info["slides"] = len(prs.slides)
        return SkillResult(success=True, data=file_info)

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def _create_pdf(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Create PDF document."""
        if HAS_FPDF:
            return self._create_pdf_binary(content, title, metadata, output_path)

        # Fallback: return markdown that can be converted
        md_content = f"# {title}\n\n{content}"
        return SkillResult(
            success=True,
            data={
                "format": "pdf",
                "markdown_source": md_content,
                "filename": self._safe_filename(title, ".pdf"),
                "note": "fpdf2 not installed — install with: pip install fpdf2"
            }
        )

    def _create_pdf_binary(self, content: str, title: str, metadata: dict, output_path: str) -> SkillResult:
        """Generate actual PDF using fpdf2."""
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Title
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(5)

        # Author / date
        if metadata.get("author") or metadata.get("date"):
            pdf.set_font("Helvetica", "I", 10)
            meta_line = ""
            if metadata.get("author"):
                meta_line += f"By {metadata['author']}"
            if metadata.get("date"):
                meta_line += f" | {metadata['date']}" if meta_line else metadata["date"]
            else:
                meta_line += f" | {datetime.now().strftime('%B %d, %Y')}" if meta_line else datetime.now().strftime('%B %d, %Y')
            pdf.cell(0, 6, meta_line, new_x="LMARGIN", new_y="NEXT", align="C")
            pdf.ln(8)

        # Content
        pdf.set_font("Helvetica", "", 11)
        lines = content.split("\n")
        for line in lines:
            stripped = line.strip()
            if not stripped:
                pdf.ln(4)
                continue
            if stripped.startswith("### "):
                pdf.set_font("Helvetica", "B", 13)
                pdf.cell(0, 8, stripped[4:], new_x="LMARGIN", new_y="NEXT")
                pdf.set_font("Helvetica", "", 11)
            elif stripped.startswith("## "):
                pdf.set_font("Helvetica", "B", 14)
                pdf.cell(0, 9, stripped[3:], new_x="LMARGIN", new_y="NEXT")
                pdf.set_font("Helvetica", "", 11)
            elif stripped.startswith("# "):
                continue  # Title already handled
            elif stripped.startswith("- ") or stripped.startswith("* "):
                pdf.cell(8, 6, chr(8226), new_x="END")
                pdf.multi_cell(0, 6, stripped[2:])
            else:
                pdf.multi_cell(0, 6, stripped)

        buffer = io.BytesIO()
        pdf.output(buffer)
        buffer.seek(0)

        file_info = self._save_or_b64(buffer, self._safe_filename(title, ".pdf"), output_path)
        file_info["format"] = "pdf"
        return SkillResult(success=True, data=file_info)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _parse_csv_content(self, content: str) -> list:
        """Parse CSV-like content into rows."""
        rows = []
        for line in content.split("\n"):
            if line.strip():
                # Simple CSV parsing (doesn't handle quoted commas)
                row = [cell.strip().strip('"') for cell in line.split(",")]
                rows.append(row)
        return rows

    def _parse_slides(self, content: str, title: str) -> list:
        """Parse content into slide structures."""
        sections = re.split(r"\n---\n", content)
        slides = []
        for i, section in enumerate(sections):
            section = section.strip()
            if not section:
                continue
            lines = section.split("\n")
            slide_title = title if i == 0 else (lines[0].lstrip("#").strip() if lines else f"Slide {i+1}")
            slides.append({
                "number": i + 1,
                "title": slide_title[:50],
                "content": section
            })
        return slides

    def _parse_markdown_to_structure(self, content: str, title: str, metadata: dict) -> dict:
        """Parse markdown content into a structured document format."""
        sections = []
        current_section = {"type": "paragraph", "content": ""}

        for line in content.split("\n"):
            stripped = line.strip()
            if stripped.startswith("### "):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"type": "heading", "level": 3, "text": stripped[4:]}
            elif stripped.startswith("## "):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"type": "heading", "level": 2, "text": stripped[3:]}
            elif stripped.startswith("# "):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"type": "heading", "level": 1, "text": stripped[2:]}
            elif stripped.startswith("- ") or stripped.startswith("* "):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"type": "paragraph", "text": stripped[2:]}
            else:
                current_section["content"] += stripped + "\n"

        if current_section.get("content") or current_section.get("text"):
            sections.append(current_section)

        return {
            "title": title,
            "author": metadata.get("author", "Aethera AI"),
            "created": datetime.now().isoformat(),
            "sections": sections
        }

    def _safe_filename(self, title: str, ext: str) -> str:
        """Convert title to a safe filename."""
        safe = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_').lower()
        return f"{safe[:50]}{ext}"